from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import os
from datetime import date

temas_colores = {
    "Clásico": {
        "mp_barras": ["#fefe00", "#92d151"],
        "ordenes_barras": ["#fefe00", "#92d151"],
        "ordenes_torta": ["#92d151", "#ffc101"],
        "mp_torta": ["#92d151", "#ffc101"]
    },
    "Corporativo": {
        "mp_barras": ["#005c5c", "#00cccc"],
        "ordenes_barras": ["#005c5c", "#00cccc"],
        "ordenes_torta": ["#339966", "#cccccc"],
        "mp_torta": ["#339966", "#cccccc"]
    }
}

def reemplazar_varias_imagenes_en_slide(slide, nuevas_imagenes):
    imagenes_reemplazadas = 0
    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # Picture
            if imagenes_reemplazadas < len(nuevas_imagenes):
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(nuevas_imagenes[imagenes_reemplazadas], left, top, width=width, height=height)
                imagenes_reemplazadas += 1

def crear_informe(nombre, mes, mp_prog, mp_ejec, ord_finalizadas, ord_pendientes, ruta_salida, tema="Clásico", plantilla="Plantilla.pptx"):
    ord_creadas = ord_finalizadas + ord_pendientes
    mp_no_ejec = max(0, mp_prog - mp_ejec)

    tema = temas_colores.get(tema, temas_colores["Clásico"])
    escala_mp = max(mp_prog, mp_ejec, 1) + 1
    escala_ordenes = max(ord_creadas, ord_finalizadas, 1) + 1

    # Gráfico 1: Barras MP
    plt.figure(figsize=(4, 3))
    valores_mp = [mp_prog, mp_ejec]
    barras = plt.bar(["Programados", "Ejecutados"], valores_mp, color=tema["mp_barras"])
    for i, barra in enumerate(barras):
        altura = barra.get_height()
        plt.text(barra.get_x() + barra.get_width()/2, altura * 0.5, str(valores_mp[i]),
                 ha='center', va='center', color='black', fontsize=10, fontweight='bold')
    plt.title("Mantenimiento Preventivo")
    plt.ylim(0, escala_mp)
    plt.yticks(range(0, escala_mp + 1, 2))
    plt.tight_layout()
    plt.savefig("grafico_mp_barras.png")
    plt.close()

    # Gráfico 2: Barras Órdenes
    plt.figure(figsize=(4, 3))
    valores_ord = [ord_creadas, ord_finalizadas]
    barras = plt.bar(["Creadas", "Finalizadas"], valores_ord, color=tema["ordenes_barras"])
    for i, barra in enumerate(barras):
        altura = barra.get_height()
        plt.text(barra.get_x() + barra.get_width()/2, altura * 0.5, str(valores_ord[i]),
                 ha='center', va='center', color='black', fontsize=10, fontweight='bold')
    plt.title("Órdenes de Mantenimiento")
    plt.ylim(0, escala_ordenes)
    plt.yticks(range(0, escala_ordenes + 1, 2))
    plt.tight_layout()
    plt.savefig("grafico_ordenes_barras.png")
    plt.close()

    # Gráfico 3: Torta Órdenes
    plt.figure(figsize=(4, 4))
    plt.pie([ord_finalizadas, ord_pendientes],
            labels=["Finalizadas", "Pendientes"],
            colors=tema["ordenes_torta"],
            autopct='%1.1f%%',
            startangle=140,
            wedgeprops=dict(edgecolor='white'))
    plt.title("Distribución de Órdenes")
    plt.axis('equal')
    plt.savefig("grafico_ordenes_torta.png")
    plt.close()

    # Gráfico 4: Torta MP
    plt.figure(figsize=(4, 4))
    plt.pie([mp_ejec, mp_no_ejec],
            labels=["Ejecutados", "No Ejecutados"],
            colors=tema["mp_torta"],
            autopct='%1.1f%%',
            startangle=140,
            wedgeprops=dict(edgecolor='white'))
    plt.title("Distribución de Mantenimientos")
    plt.axis('equal')
    plt.savefig("grafico_mp_torta.png")
    plt.close()

    prs = Presentation(plantilla)
    nombre = nombre.upper()
    mes = mes.upper()

    for i in [0, 1]:
        for shape in prs.slides[i].shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                nuevo_texto = text.replace("INGENIERO", nombre).replace("MES", mes)
                shape.text_frame.text = nuevo_texto

    if len(prs.slides) >= 5:
        reemplazar_varias_imagenes_en_slide(prs.slides[3], ["grafico_mp_barras.png", "grafico_mp_torta.png"])
        reemplazar_varias_imagenes_en_slide(prs.slides[4], ["grafico_ordenes_barras.png", "grafico_ordenes_torta.png"])

    prs.save(ruta_salida)

    # Limpiar
    for nombre in [
        "grafico_mp_barras.png",
        "grafico_mp_torta.png",
        "grafico_ordenes_barras.png",
        "grafico_ordenes_torta.png"
    ]:
        try:
            os.remove(nombre)
        except:
            pass
